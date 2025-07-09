
import streamlit as st
import pandas as pd
import numpy as np
from io import BytesIO
import time
import requests
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

# Seed for reproducibility
np.random.seed(42)

# Page config
st.set_page_config(page_title="MerchSentAI", page_icon="🛡️", layout="centered")

# Initialize session state
if 'page' not in st.session_state:
    st.session_state.page = 'splash'

# Helper: compute flags and scores
def add_flags(df, threshold=50):
    df = df.copy()
    df['Enriched_RiskScore'] = df.get('Enriched_RiskScore', 0)
    # Define rules
    results = {
        'high_risk_countries': ['Iran','North Korea','Syria','Cuba','Russia','Venezuela'],
        'high_risk_categories': ['Cryptocurrency','Gambling','Adult Content','Arms Dealing'],
        'pep_list': ['John Smith','Emma Johnson','Michael Brown','Sarah Davis']
    }
    df['RegistrationDate'] = pd.to_datetime(df['RegistrationDate'].str.strip(), errors='coerce')
    rules = {
        'SanctionedFlag': df['MerchantName'].isin(results['pep_list']),
        'HighRiskCountryFlag': df['Country'].isin(results['high_risk_countries']),
        'HighRiskCategoryFlag': df['BusinessCategory'].isin(results['high_risk_categories']),
        'HighAmountFlag': df['AvgTransaction'] > 10000,
        'RecentRegistrationFlag': (pd.Timestamp('today') - df['RegistrationDate']).dt.days < 30,
        'CryptoRelatedFlag': df['BusinessCategory']=='Cryptocurrency',
        'LuxuryHighValueFlag': (df['BusinessCategory']=='Luxury Goods') & (df['AvgTransaction']>50000),
        'EnergyLargeVolumeFlag': (df['BusinessCategory']=='Energy') & (df['AvgTransaction']>100000),
        'OddNameLengthFlag': df['MerchantName'].str.len() % 2 == 1,
        'HighRiskScoreFlag': df['Enriched_RiskScore'] > 80
    }
    for name, cond in rules.items():
        df[name] = cond.astype(int)
    flag_cols = list(rules.keys())
    # Calculate FraudScore and Suspicious
    df['FlagsTripped'] = df[flag_cols].sum(axis=1)
    df['FraudScore'] = (df['FlagsTripped'] / len(flag_cols) * 100).round(1)
    df['IsSuspicious'] = df['FraudScore'] >= threshold
    return df, flag_cols

# Cost simulation
def simulate_cost(df, fraud_cost, sanction_cost):
    sanction_misses = int(df['SanctionedFlag'].sum())
    fraud_misses = int(df['IsSuspicious'].sum())
    sanction_total = sanction_misses * sanction_cost
    fraud_total = fraud_misses * fraud_cost
    total = sanction_total + fraud_total
    return sanction_misses, sanction_total, fraud_misses, fraud_total, total

# Splash
def splash_page():
    st.markdown("""
    <style>
    .stApp {background: linear-gradient(135deg,#e8f0ff,#f3e8ff);}
    .stButton>button{background-color:#4f46e5;color:white;font-weight:bold;
      border-radius:8px;padding:.75rem;width:100%;}
    </style>
    """, unsafe_allow_html=True)
    logo = st.file_uploader("Upload Logo", type=["png","jpg","jpeg"], key="logo")
    if logo:
        st.image(logo, use_column_width=True)
    st.title("Merchant Risk Intelligence Platform")
    st.write("Know your merchant, know your risk.")
    if st.button("Enter the Data World"):
        st.session_state.page = 'analysis'

# Main analysis
def analysis_page():
    st.title("Merchant Risk Analysis")
    upload = st.file_uploader("Upload Merchant Data (CSV/XLS/XLSX)", type=["csv","xls","xlsx"], key="data")
    if not upload:
        st.info("Please upload data to proceed.")
        return

    # Load on new file
    if ('raw_df' not in st.session_state) or (st.session_state.filename != upload.name):
        try:
            df = pd.read_csv(upload) if upload.name.lower().endswith('.csv') else pd.read_excel(upload)
            df.columns = df.columns.str.strip()
            st.session_state.raw_df = df.copy()
            st.session_state.df = df.copy()
            st.session_state.filename = upload.name
            for flag in ['enriched','screened','fraud_done','processor_done']:
                st.session_state[flag] = False
        except Exception as e:
            st.error(f"Error reading data: {e}")
            return

    df = st.session_state.df

    # Stage 1
    st.subheader("Stage 1: Data Enrichment")
    if not st.session_state.enriched and st.button("Enrich Data"):
        df['Enriched_Location'] = df.get('Country', np.nan)
        df['Enriched_RiskScore'] = np.random.randint(1,101,len(df))
        st.session_state.df = df
        st.session_state.enriched = True
        st.success("Data enriched")
    if st.session_state.enriched:
        st.table(df.head())

    # Stage 2
    st.subheader("Stage 2: Perform Data Sanctions Screening")
    if st.session_state.enriched and not st.session_state.screened and st.button("Perform Data Sanctions Screening"):
        peps = ['John Smith','Emma Johnson','Michael Brown','Sarah Davis']
        df['SanctionedFlag'] = df['MerchantName'].isin(peps).astype(int)
        st.session_state.df = df
        st.session_state.screened = True
        st.success("Screening complete")
    if st.session_state.screened:
        checks = pd.DataFrame({'Check':['Sanction Hits'],'Hits':[int(df['SanctionedFlag'].sum())]})
        st.table(checks)

    # Stage 3
    st.subheader("Stage 3: Run Fraud Analysis")
    threshold = st.slider("Suspicious Threshold (%)", 0, 100, 50, step=5)
    if st.session_state.screened and not st.session_state.fraud_done and st.button("Run Fraud Analysis"):
        df, flag_cols = add_flags(df, threshold)
        st.session_state.df = df
        st.session_state.flag_cols = flag_cols
        st.session_state.fraud_done = True
        st.success("Fraud analysis complete")
    if st.session_state.fraud_done:
        summary = pd.DataFrame({
            'Metric':['Total Merchants','Suspicious Merchants'],
            'Count':[len(df), int(df['IsSuspicious'].sum())]
        })
        st.table(summary)
        susp = df[df['IsSuspicious']][['MerchantName','FraudScore','FlagsTripped']+st.session_state.flag_cols]
        st.subheader("Suspicious Merchants Detail")
        st.table(susp)
        buf = BytesIO()
        with pd.ExcelWriter(buf, engine='xlsxwriter') as writer:
            summary.to_excel(writer, sheet_name='FraudSummary', index=False)
            susp.to_excel(writer, sheet_name='SuspiciousList', index=False)
        buf.seek(0)
        st.download_button("Download Fraud Summary (Excel)", buf.getvalue(),
                           file_name="fraud_summary.xlsx",
                           mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

    # Stage 4
    st.subheader("Stage 4: Payment Processor Impact")
    fraud_cost = st.number_input("Cost per fraud incident ($)", min_value=0, value=200, step=50)
    sanction_cost = st.number_input("Cost per sanction hit ($)", min_value=0, value=5000, step=500)
    if st.session_state.fraud_done and not st.session_state.processor_done and st.button("Simulate Processor Impact"):
        # simulate API call
        with st.spinner("Calling MerchSentAI API..."):
            start = time.time()
            time.sleep(0.2)
            latency = (time.time() - start)*1000
        st.success(f"API response OK in {latency:.0f} ms")
        # simulate costs
        raw_df = st.session_state.raw_df.copy()
        raw_df, _ = add_flags(raw_df, threshold)
        rm = simulate_cost(raw_df, fraud_cost, sanction_cost)
        ref = simulate_cost(df, fraud_cost, sanction_cost)
        st.session_state.raw_metrics = rm
        st.session_state.ref_metrics = ref
        st.session_state.processor_done = True

    if st.session_state.processor_done:
        rm = st.session_state.raw_metrics
        ref = st.session_state.ref_metrics
        # display infographic
        img = Image.open("/mnt/data/A_2D_digital_infographic_compares_two_approaches_t.png")
        st.image(img, use_column_width=True, caption="Before vs. After: Processor Impact")
        # KPIs side-by-side
        col1, col2 = st.columns(2)
        with col1:
            st.metric("Raw Processor Loss", f"${rm[4]:,}", delta=f"{rm[2]} incidents")
        with col2:
            st.metric("MerchSentAI Processor Loss", f"${ref[4]:,}", delta=f"{ref[2]} incidents")
        savings = rm[4] - ref[4]
        pct = savings/rm[4]*100 if rm[4] else 0
        st.markdown(f"### 💰 Savings: **${savings:,}** ({pct:.1f}% reduction)")
        # Problem merchants
        problem = st.session_state.raw_df.copy()
        problem, _ = add_flags(problem, threshold)
        problem = problem[problem['IsSuspicious']][['MerchantName','FraudScore','FlagsTripped']]
        st.subheader("Problem Merchants (Raw Data)")
        st.table(problem)
        # Download impact report
        buf2 = BytesIO()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Processor Impact Summary"
        slide.placeholders[1].text = f"Savings: ${savings:,} ({pct:.1f}% reduction)"
        prs.save(buf2)
        buf2.seek(0)
        st.download_button("Download Impact PPTX", buf2.getvalue(),
                           file_name="impact_summary.pptx",
                           mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    if st.button("Back to Splash"):
        st.session_state.page = 'splash'

# Run
if st.session_state.page == 'splash':
    splash_page()
else:
    analysis_page()
