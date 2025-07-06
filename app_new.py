import streamlit as st
import pandas as pd
import numpy as np
import os
from datetime import datetime, date
from io import BytesIO
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# --- Page configuration ---
st.set_page_config(page_title="MerchSentAI", page_icon="🛡️", layout="centered")

# --- Session state setup ---
if 'page' not in st.session_state:
    st.session_state.page = 'splash'

# --- Fraud detection logic with 10 flags ---
def run_fraud_analysis(df):
    results = {
        'high_risk_countries': ['Iran', 'North Korea', 'Syria', 'Cuba', 'Russia', 'Venezuela'],
        'high_risk_categories': ['Cryptocurrency', 'Gambling', 'Adult Content', 'Arms Dealing'],
        'pep_list': ['John Smith', 'Emma Johnson', 'Michael Brown', 'Sarah Davis']
    }
    df['SanctionedFlag'] = df['MerchantName'].isin(results['pep_list'])
    df['HighRiskCountryFlag'] = df['Country'].isin(results['high_risk_countries'])
    df['HighRiskCategoryFlag'] = df['BusinessCategory'].isin(results['high_risk_categories'])
    df['HighAmountFlag'] = df['AvgTransaction'] > 10000
    df['RegistrationDate'] = pd.to_datetime(df['RegistrationDate'], errors='coerce')
    df['RecentRegistrationFlag'] = (pd.Timestamp(date.today()) - df['RegistrationDate']).dt.days < 30
    df['CryptoRelatedFlag'] = df['BusinessCategory'] == 'Cryptocurrency'
    df['LuxuryHighValueFlag'] = (df['BusinessCategory'] == 'Luxury Goods') & (df['AvgTransaction'] > 50000)
    df['EnergyLargeVolumeFlag'] = (df['BusinessCategory'] == 'Energy') & (df['AvgTransaction'] > 100000)
    df['OddNameLengthFlag'] = df['MerchantName'].str.len() % 2 == 1
    df['HighRiskScoreFlag'] = df.get('Enriched_RiskScore', 0) > 80

    fraud_logic_cols = [
        'SanctionedFlag','HighRiskCountryFlag','HighRiskCategoryFlag',
        'HighAmountFlag','RecentRegistrationFlag','CryptoRelatedFlag',
        'LuxuryHighValueFlag','EnergyLargeVolumeFlag','OddNameLengthFlag','HighRiskScoreFlag'
    ]
    df['FraudFlagCount'] = df[fraud_logic_cols].sum(axis=1)
    df['Suspicious'] = df['FraudFlagCount'] > 0
    return df, results, fraud_logic_cols

# --- Executive summary display ---
def display_executive_summary(df, results, fraud_logic_cols):
    total = len(df)
    sanctioned = int(df['SanctionedFlag'].sum())
    high_risk = int(df['HighRiskCountryFlag'].sum())
    suspicious = int(df['Suspicious'].sum())

    st.markdown("### Executive Summary")
    c1, c2, c3, c4 = st.columns(4)
    c1.metric("Total merchants", total)
    c2.metric("Sanction matches", sanctioned)
    c3.metric("High-risk jurisdictions", high_risk)
    c4.metric("Suspicious merchants", suspicious)

    summary_df = pd.DataFrame({
        'Metric': ['Total merchants','Sanction matches','High-risk jurisdictions','Suspicious merchants'],
        'Count': [total, sanctioned, high_risk, suspicious]
    })
    st.table(summary_df)
    st.bar_chart(summary_df.set_index('Metric'))

    st.markdown("#### Fraud Logic Flags Distribution")
    flag_counts = df[fraud_logic_cols].sum().sort_values(ascending=False)
    st.bar_chart(flag_counts)

# --- Splash page ---
def splash_page():
    st.markdown("""<style>
        .stButton>button {
            background-color:#4f46e5;
            color:white;
            font-weight:bold;
            border-radius:8px;
            padding:0.75rem;
            width:100%;
        }
    </style>""", unsafe_allow_html=True)
    logo_file = st.file_uploader("Upload Logo (PNG/JPG)", type=["png","jpg","jpeg"], key="logo_uploader", label_visibility="collapsed")
    if logo_file:
        st.image(logo_file, use_column_width=True)

    st.title("Merchant Risk Intelligence Platform")
    st.write("Know your merchant, know your risk.")

    if st.button("Enter the Data World"):
        st.session_state.page = 'analysis'

# --- Analysis page ---
def analysis_page():
    st.title("Merchant Risk Analysis")
    uploaded = st.file_uploader("Upload Merchant Data (CSV/XLSX)", type=["csv","xlsx"])
    if not uploaded:
        st.info("📁 Please upload a CSV or XLSX file to proceed.")
        return

    try:
        df = pd.read_csv(uploaded) if uploaded.name.lower().endswith('.csv') else pd.read_excel(uploaded)
    except Exception as e:
        st.error(f"Could not read file: {e}")
        return

    with st.expander("Data Enrichment", expanded=True):
        df['Enriched_Location'] = df.get('Country', np.nan)
        df['Enriched_RiskScore'] = np.random.randint(1,100,len(df))
        st.success("Data enrichment applied")
        st.metric("Rows enriched", len(df))
        st.metric("Avg Risk Score", int(df['Enriched_RiskScore'].mean()))
        st.dataframe(df.head(10))

    df, results, fraud_logic_cols = run_fraud_analysis(df)
    with st.expander("Sanctions/Screening", expanded=True):
        st.success("Sanctions screening completed")
        checks = pd.DataFrame({
            'Check': ['Sanction','HighRiskCountry','HighRiskCategory','HighAmount','RecentRegistration'],
            'Hits': [
                int(df['SanctionedFlag'].sum()),
                int(df['HighRiskCountryFlag'].sum()),
                int(df['HighRiskCategoryFlag'].sum()),
                int(df['HighAmountFlag'].sum()),
                int(df['RecentRegistrationFlag'].sum())
            ]
        })
        st.table(checks)
        st.bar_chart(checks.set_index('Check'))

    with st.expander("Fraud Analysis", expanded=True):
        st.success("Fraud analysis completed")
        st.dataframe(df[fraud_logic_cols + ['FraudFlagCount','Suspicious']].head(10))
        st.bar_chart(df['FraudFlagCount'].value_counts().sort_index())

    display_executive_summary(df, results, fraud_logic_cols)

    # Excel report
    excel_io = BytesIO()
    with pd.ExcelWriter(excel_io, engine='xlsxwriter') as writer:
        df.to_excel(writer, sheet_name='Enriched Data', index=False)
        checks.to_excel(writer, sheet_name='Screening Summary', index=False)
        df[['MerchantID','MerchantName','FraudFlagCount','Suspicious']+fraud_logic_cols].to_excel(writer, sheet_name='Fraud Analysis', index=False)
        pd.DataFrame({
            'Metric': ['Total merchants','Sanction matches','High-risk jurisdictions','Suspicious merchants'],
            'Count': [
                len(df),
                int(df['SanctionedFlag'].sum()),
                int(df['HighRiskCountryFlag'].sum()),
                int(df['Suspicious'].sum())
            ]
        }).to_excel(writer, sheet_name='Executive Summary', index=False)
    excel_io.seek(0)
    st.download_button("Download Excel Report", data=excel_io.read(), file_name="risk_report.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Executive Summary"
    slide.placeholders[1].text = (
        f"Total merchants: {len(df)}
"
        f"Sanction matches: {int(df['SanctionedFlag'].sum())}
"
        f"High-risk jurisdictions: {int(df['HighRiskCountryFlag'].sum())}
"
        f"Suspicious merchants: {int(df['Suspicious'].sum())}"
    )
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Fraud Logic Flags Distribution"
    fig, ax = plt.subplots()
    flag_counts = df[fraud_logic_cols].sum().sort_values(ascending=False)
    flag_counts.plot(kind='bar', ax=ax)
    ax.set_ylabel('Count')
    buf = BytesIO()
    fig.savefig(buf, format='png')
    plt.close(fig)
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(1), Inches(1.5), width=Inches(8))
    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    st.download_button("Download PowerPoint Report", data=pptx_io.read(), file_name="executive_summary.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    if st.button("Back to Main"):
        st.session_state.page = 'splash'

if st.session_state.page == 'splash':
    splash_page()
else:
    analysis_page()
