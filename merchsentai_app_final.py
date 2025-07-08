
import streamlit as st
import pandas as pd
import numpy as np
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches

# Seed for reproducibility
np.random.seed(42)

# Page config
st.set_page_config(page_title="MerchSentAI", page_icon="🛡️", layout="centered")

# Initialize session state
if 'page' not in st.session_state:
    st.session_state.page = 'splash'

# Flag computation helper
def add_flags(df):
    df = df.copy()
    df['Enriched_RiskScore'] = df.get('Enriched_RiskScore', 0)
    results = {
        'high_risk_countries': ['Iran', 'North Korea', 'Syria', 'Cuba', 'Russia', 'Venezuela'],
        'high_risk_categories': ['Cryptocurrency', 'Gambling', 'Adult Content', 'Arms Dealing'],
        'pep_list': ['John Smith', 'Emma Johnson', 'Michael Brown', 'Sarah Davis']
    }
    df['RegistrationDate'] = pd.to_datetime(df['RegistrationDate'].str.strip(), errors='coerce')

    # Fraud rules
    rules = {
        'SanctionedFlag': df['MerchantName'].isin(results['pep_list']),
        'HighRiskCountryFlag': df['Country'].isin(results['high_risk_countries']),
        'HighRiskCategoryFlag': df['BusinessCategory'].isin(results['high_risk_categories']),
        'HighAmountFlag': df['AvgTransaction'] > 10000,
        'RecentRegistrationFlag': (pd.Timestamp('today') - df['RegistrationDate']).dt.days < 30,
        'CryptoRelatedFlag': df['BusinessCategory']=='Cryptocurrency',
        'LuxuryHighValueFlag': (df['BusinessCategory']=='Luxury Goods') & (df['AvgTransaction']>50000),
        'EnergyLargeVolumeFlag': (df['BusinessCategory']=='Energy') & (df['AvgTransaction']>100000),
        'OddNameLengthFlag': df['MerchantName'].str.len() % 2 == 1,
        'HighRiskScoreFlag': df['Enriched_RiskScore'] > 80
    }
    for name, condition in rules.items():
        df[name] = condition.astype(int)
    flag_cols = list(rules.keys())

    # Compute fraud score: percent of rules triggered
    df['FlagsTripped'] = df[flag_cols].sum(axis=1)
    df['FraudScore'] = (df['FlagsTripped'] / len(flag_cols) * 100).round(1)
    df['IsSuspicious'] = df['FraudScore'] >= 50

    return df, flag_cols

# Cost simulation
def simulate_cost(df, flag_cols, fraud_cost, sanction_cost):
    sanction_misses = int(df['SanctionedFlag'].sum())
    fraud_misses = int(df['IsSuspicious'].sum())
    sanction_total = sanction_misses * sanction_cost
    fraud_total = fraud_misses * fraud_cost
    total = sanction_total + fraud_total
    return sanction_misses, sanction_total, fraud_misses, fraud_total, total

# Splash page
def splash_page():
    st.markdown("""
    <style>
    .stApp {background: linear-gradient(135deg, #e8f0ff, #f3e8ff);}
    .stButton>button {background-color:#4f46e5; color:white; font-weight:bold; 
                       border-radius:8px; padding:.75rem; width:100%;}
    </style>
    """, unsafe_allow_html=True)
    logo = st.file_uploader("Upload Logo (PNG/JPG)", type=["png","jpg","jpeg"], key="logo")
    if logo:
        st.image(logo, use_column_width=True)
    st.title("Merchant Risk Intelligence Platform")
    st.write("Know your merchant, know your risk.")
    if st.button("Enter the Data World"):
        st.session_state.page = 'analysis'

# Analysis page
def analysis_page():
    st.title("Merchant Risk Analysis")
    upload = st.file_uploader("Upload Merchant Data (CSV/XLS/XLSX)", 
                              type=["csv","xls","xlsx"], key="data")
    if not upload:
        st.info("Please upload data to proceed.")
        return

    # Load data
    if ('raw_df' not in st.session_state) or (st.session_state.filename != upload.name):
        try:
            name = upload.name.lower()
            df = pd.read_csv(upload) if name.endswith('.csv') else pd.read_excel(upload)
            df.columns = df.columns.str.strip()
            st.session_state.raw_df = df.copy()
            st.session_state.df = df.copy()
            st.session_state.filename = upload.name
            st.session_state.enriched = False
            st.session_state.screened = False
            st.session_state.fraud_done = False
            st.session_state.processor_done = False
        except Exception as e:
            st.error(f"Failed to read file: {e}")
            return

    df = st.session_state.df

    # Stage 1: Enrichment
    st.subheader("Stage 1: Data Enrichment")
    if not st.session_state.enriched and st.button("Enrich Data"):
        df['Enriched_Location'] = df.get('Country', np.nan)
        df['Enriched_RiskScore'] = np.random.randint(1, 101, len(df))
        st.session_state.df = df
        st.session_state.enriched = True
        st.success("Data enriched!")
    if st.session_state.enriched:
        st.table(df.head())

    # Stage 2: Sanctions Screening
    st.subheader("Stage 2: Perform Data Sanctions Screening")
    if st.session_state.enriched and not st.session_state.screened and st.button("Perform Data Sanctions Screening"):
        peps = ['John Smith','Emma Johnson','Michael Brown','Sarah Davis']
        df['SanctionedFlag'] = df['MerchantName'].isin(peps).astype(int)
        st.session_state.df = df
        st.session_state.screened = True
        st.success("Sanctions screening complete!")
    if st.session_state.screened:
        checks = pd.DataFrame({
            'Check': ['Sanction hits'],
            'Hits': [int(df['SanctionedFlag'].sum())]
        })
        st.table(checks)

    # Stage 3: Fraud Analysis
    st.subheader("Stage 3: Run Fraud Analysis")
    if st.session_state.screened and not st.session_state.fraud_done and st.button("Run Fraud Analysis"):
        df, flag_cols = add_flags(df)
        st.session_state.df = df
        st.session_state.flag_cols = flag_cols
        st.session_state.fraud_done = True
        st.success("Fraud analysis complete!")
    if st.session_state.fraud_done:
        # Summary table
        summary = pd.DataFrame({
            'Metric': ['Total merchants', 'Suspicious merchants'],
            'Count': [len(df), int(df['IsSuspicious'].sum())]
        })
        st.table(summary)
        # Suspicious merchants detail
        suspicious = df[df['IsSuspicious']][['MerchantName','FraudScore','FlagsTripped'] + st.session_state.flag_cols]
        st.subheader("Suspicious Merchants Detail")
        st.table(suspicious)

        # Executive summary download
        buf = BytesIO()
        with pd.ExcelWriter(buf, engine='xlsxwriter') as writer:
            summary.to_excel(writer, sheet_name='Fraud Summary', index=False)
            suspicious.to_excel(writer, sheet_name='Suspicious Merchants', index=False)
        buf.seek(0)
        st.download_button("Download Fraud Executive Summary (Excel)", 
                           buf.getvalue(), 
                           file_name="fraud_summary.xlsx",
                           mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

    # Stage 4: Payment Processor Impact
    st.subheader("Stage 4: Payment Processor Impact")
    fraud_cost = st.number_input("Cost per fraud incident ($)", min_value=0, value=200, step=50)
    sanction_cost = st.number_input("Cost per sanction hit ($)", min_value=0, value=5000, step=500)
    if st.session_state.fraud_done and not st.session_state.processor_done and st.button("Simulate Processor Impact"):
        # simulate costs using raw and refined
        raw_flags_df, _ = add_flags(st.session_state.raw_df)
        raw_metrics = simulate_cost(raw_flags_df, st.session_state.flag_cols, fraud_cost, sanction_cost)
        ref_metrics = simulate_cost(df, st.session_state.flag_cols, fraud_cost, sanction_cost)
        st.session_state.raw_metrics = raw_metrics
        st.session_state.ref_metrics = ref_metrics
        st.session_state.processor_done = True
        st.success("Processor impact simulated!")
    if st.session_state.processor_done:
        rm = st.session_state.raw_metrics
        fm = st.session_state.ref_metrics
        impact = pd.DataFrame({
            'Issue': ['Sanction Hits','Suspicious Incidents','Total Loss'],
            'Raw Count/Cost': [f"{rm[0]} / ${rm[1]}", f"{rm[2]} / ${rm[3]}", f"${rm[4]}"],
            'Refined Count/Cost': [f"{fm[0]} / ${fm[1]}", f"{fm[2]} / ${fm[3]}", f"${fm[4]}"],
            'Savings': ['', '', f"${rm[4]-fm[4]}"]
        })
        st.table(impact)
        # Summary text
        total_savings = rm[4] - fm[4]
        total_raw = rm[4]
        pct = (total_savings/total_raw*100) if total_raw>0 else 0
        st.markdown(f"**MerchSentAI Impact Summary:** By enriching, screening, and fraud reviewing data, MerchSentAI helped the payment processor avoid **${{:,}}** in potential losses, a reduction of {pct:.1f}%.".format(total_savings))

        # Download full impact report
        buf2 = BytesIO()
        prs = Presentation()
        sl = prs.slides.add_slide(prs.slide_layouts[0])
        sl.shapes.title.text = "Processor Impact Summary"
        text = sl.shapes.placeholders[1].text_frame
        text.text = f"Total Loss Avoided: ${total_savings}"
        prs.save(buf2)
        buf2.seek(0)
        st.download_button("Download Impact Report (PPTX)", 
                           buf2.getvalue(), 
                           file_name="impact_summary.pptx",
                           mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    # Back to splash
    if st.button("Back to Splash"):
        st.session_state.page = 'splash'

# Run proper page
if st.session_state.page == 'splash':
    splash_page()
else:
    analysis_page()
