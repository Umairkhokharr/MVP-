
import streamlit as st
import pandas as pd
import numpy as np
from io import BytesIO
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# Seed for reproducibility in demos
np.random.seed(42)

# --- Page configuration ---
st.set_page_config(page_title="MerchSentAI", page_icon="🛡️", layout="centered")

# --- Session state initialization ---
if 'page' not in st.session_state:
    st.session_state.page = 'splash'

# --- Fraud analysis helper ---
def run_fraud_analysis(df):
    # Ensure enrichment score exists
    if 'Enriched_RiskScore' not in df:
        st.warning("Enriched_RiskScore missing—defaulting to 0")
        df['Enriched_RiskScore'] = 0

    # Convert RegistrationDate to datetime
    df['RegistrationDate'] = pd.to_datetime(df['RegistrationDate'].str.strip(), errors='coerce')

    # Fraud logic flags
    df['CryptoFlag'] = df['BusinessCategory'] == 'Cryptocurrency'
    df['LuxuryFlag'] = (df['BusinessCategory']=='Luxury Goods') & (df['AvgTransaction']>50000)
    df['OddNameFlag'] = df['MerchantName'].str.len() % 2 == 1
    df['HighScoreFlag'] = df['Enriched_RiskScore'] > 80

    flags = ['CryptoFlag','LuxuryFlag','OddNameFlag','HighScoreFlag']
    # Summary
    summary = pd.DataFrame({
        'Metric': ['Total merchants','Flags tripped'],
        'Count': [len(df), int(df[flags].any(axis=1).sum())]
    })
    return df, summary, flags

# --- Splash page ---
def splash_page():
    st.markdown("""
    <style>
    .stApp {background: linear-gradient(135deg, #e8f0ff, #f3e8ff);}
    .stButton>button { background-color:#4f46e5; color:white; font-weight:bold; border-radius:8px; padding:.75rem; width:100%; }
    </style>
    """, unsafe_allow_html=True)
    logo = st.file_uploader("Upload Logo", type=["png","jpg","jpeg"], key="logo")
    if logo:
        st.image(logo, use_column_width=True)
    st.title("Merchant Risk Intelligence")
    st.write("Know your merchant, know your risk.")
    if st.button("Enter the Data World"):
        st.session_state.page = 'analysis'

# --- Analysis page with staged buttons ---
def analysis_page():
    st.title("Merchant Risk Analysis")
    upload = st.file_uploader("Upload Merchant Data (CSV/XLS/XLSX)", type=["csv","xls","xlsx"], key="data")
    if upload:
        if ('df' not in st.session_state) or (st.session_state.get('filename')!=upload.name):
            # Load and reset
            try:
                name = upload.name.lower()
                if name.endswith('.csv'):
                    df = pd.read_csv(upload)
                else:
                    df = pd.read_excel(upload)
                df.columns = df.columns.str.strip()
                st.session_state.df = df
                st.session_state.enriched = False
                st.session_state.screened = False
                st.session_state.fraud_done = False
                st.session_state.filename = upload.name
            except Exception as e:
                st.error(f"Failed to read file: {e}")
                return
    else:
        st.info("Please upload data to proceed.")
        return

    df = st.session_state.df

    # Stage 1: Enrichment
    st.subheader("Stage 1: Data Enrichment")
    if not st.session_state.enriched:
        if st.button("Enrich Data"):
            df['Enriched_Location'] = df.get('Country', np.nan)
            df['Enriched_RiskScore'] = np.random.randint(1,100,len(df))
            st.session_state.enriched = True
            st.success("Data enriched!")
    if st.session_state.enriched:
        st.dataframe(df.head())

    # Stage 2: Sanctions Screening
    st.subheader("Stage 2: Perform Data Sanctions Screening")
    if st.session_state.enriched and not st.session_state.screened:
        if st.button("Perform Data Sanctions Screening"):
            # Simple screening logic
            peps = ['John Smith','Emma Johnson','Michael Brown','Sarah Davis']
            df['SanctionedFlag'] = df['MerchantName'].isin(peps)
            st.session_state.screened = True
            st.success("Sanctions screening complete!")
    if st.session_state.screened:
        checks = pd.DataFrame({
            'Check': ['Sanction hits'],
            'Hits': [int(st.session_state.df['SanctionedFlag'].sum())]
        })
        st.table(checks)

    # Stage 3: Fraud Analysis
    st.subheader("Stage 3: Run Fraud Analysis")
    if st.session_state.screened and not st.session_state.fraud_done:
        if st.button("Run Fraud Analysis"):
            df, summary, flags = run_fraud_analysis(df)
            st.session_state.summary = summary
            st.session_state.flags = flags
            st.session_state.fraud_done = True
            st.success("Fraud analysis finished!")
    if st.session_state.fraud_done:
        st.markdown("### Fraud Summary")
        st.table(st.session_state.summary)
        st.bar_chart(st.session_state.summary.set_index('Metric'))
        flag_counts = df[st.session_state.flags].sum().sort_values(ascending=False)
        st.markdown("### Fraud Flags Distribution")
        st.bar_chart(flag_counts)

        # Executive summary downloads
        buf = BytesIO()
        with pd.ExcelWriter(buf, engine='xlsxwriter') as writer:
            st.session_state.summary.to_excel(writer, index=False, sheet_name='Summary')
            writer.close()
        buf.seek(0)
        st.download_button("Download Excel Report", data=buf.getvalue(),
                           file_name="executive_summary.xlsx",
                           mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Executive Summary"
        slide.placeholders[1].text = "\n".join(
            f"{r['Metric']}: {r['Count']}" for _, r in st.session_state.summary.iterrows()
        )
        buf2 = BytesIO()
        fig, ax = plt.subplots(figsize=(6,4))
        flag_counts.plot(kind='bar', ax=ax)
        plt.tight_layout()
        fig.savefig(buf2, format='png')
        buf2.seek(0)
        slide2 = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts)>5 else prs.slide_layouts[0])
        slide2.shapes.title.text = "Fraud Flags"
        slide2.shapes.add_picture(buf2, Inches(1), Inches(1.5), width=Inches(8))

        pptx_buf = BytesIO()
        prs.save(pptx_buf)
        pptx_buf.seek(0)
        st.download_button("Download PowerPoint Report", data=pptx_buf.getvalue(),
                           file_name="executive_summary.pptx",
                           mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    # Back button
    if st.button("Back to Splash"):
        st.session_state.page = 'splash'

# --- Run pages ---
if st.session_state.page == 'splash':
    splash_page()
else:
    analysis_page()
