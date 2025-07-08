import streamlit as st
import pandas as pd
import numpy as np
from io import BytesIO
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# Seed for reproducibility in demos
np.random.seed(42)

# --- Page configuration ---
st.set_page_config(page_title="MerchSentAI", page_icon="🛡️", layout="centered")

# --- Session state setup ---
if 'page' not in st.session_state:
    st.session_state.page = 'splash'

# --- Fraud analysis logic ---
def run_fraud_analysis(df):
    # Ensure enrichment score exists
    if 'Enriched_RiskScore' not in df:
        st.warning("Enriched_RiskScore column missing—defaulting to 0")
        df['Enriched_RiskScore'] = 0

    results = {
        'high_risk_countries': ['Iran', 'North Korea', 'Syria', 'Cuba', 'Russia', 'Venezuela'],
        'high_risk_categories': ['Cryptocurrency', 'Gambling', 'Adult Content', 'Arms Dealing'],
        'pep_list': ['John Smith', 'Emma Johnson', 'Michael Brown', 'Sarah Davis']
    }

    # Convert RegistrationDate to datetime
    df['RegistrationDate'] = pd.to_datetime(df['RegistrationDate'].str.strip(), errors='coerce')

    # Core flags
    df['SanctionedFlag'] = df['MerchantName'].isin(results['pep_list'])
    df['HighRiskCountryFlag'] = df['Country'].isin(results['high_risk_countries'])
    df['HighRiskCategoryFlag'] = df['BusinessCategory'].isin(results['high_risk_categories'])
    df['HighAmountFlag'] = df['AvgTransaction'] > 10000
    df['RecentRegistrationFlag'] = (pd.Timestamp('today') - df['RegistrationDate']).dt.days < 30

    # Extended flags
    df['CryptoRelatedFlag'] = df['BusinessCategory'] == 'Cryptocurrency'
    df['LuxuryHighValueFlag'] = (df['BusinessCategory'] == 'Luxury Goods') & (df['AvgTransaction'] > 50000)
    df['EnergyLargeVolumeFlag'] = (df['BusinessCategory'] == 'Energy') & (df['AvgTransaction'] > 100000)
    df['OddNameLengthFlag'] = df['MerchantName'].str.len() % 2 == 1
    df['HighRiskScoreFlag'] = df['Enriched_RiskScore'] > 80

    fraud_logic_cols = [
        'SanctionedFlag', 'HighRiskCountryFlag', 'HighRiskCategoryFlag',
        'HighAmountFlag', 'RecentRegistrationFlag', 'CryptoRelatedFlag',
        'LuxuryHighValueFlag', 'EnergyLargeVolumeFlag', 'OddNameLengthFlag', 'HighRiskScoreFlag'
    ]

    # Summary table
    total = len(df)
    sanctioned = int(df['SanctionedFlag'].sum())
    high_risk = int(df['HighRiskCountryFlag'].sum())
    suspicious = int(df[fraud_logic_cols].any(axis=1).sum())
    summary_df = pd.DataFrame({
        'Metric': ['Total merchants', 'Sanction matches', 'High-risk jurisdictions', 'Suspicious merchants'],
        'Count': [total, sanctioned, high_risk, suspicious]
    })
    st.table(summary_df)
    st.bar_chart(summary_df.set_index('Metric'))

    st.markdown("#### Fraud Logic Flags Distribution")
    flag_counts = df[fraud_logic_cols].sum().sort_values(ascending=False)
    st.bar_chart(flag_counts)

    return df, summary_df, fraud_logic_cols

# --- Splash page ---
def splash_page():
    st.markdown(
        """
        <style>
        .stApp {background: linear-gradient(135deg, #e8f0ff, #f3e8ff);}
        .stButton>button {
            background-color: #4f46e5;
            color: white;
            font-weight: bold;
            border-radius: 8px;
            padding: 0.75rem;
            width: 100%;
        }
        </style>
        """,
        unsafe_allow_html=True
    )
    logo_file = st.file_uploader(
        "Upload Logo (PNG/JPG)", type=["png", "jpg", "jpeg"], key="logo_uploader", label_visibility="collapsed"
    )
    if logo_file:
        st.image(logo_file, use_column_width=True)

    st.title("Merchant Risk Intelligence Platform")
    st.write("Know your merchant, know your risk.")

    if st.button("Enter the Data World"):
        st.session_state.page = 'analysis'

# --- Analysis page ---
def analysis_page():
    st.title("Merchant Risk Analysis")
    uploaded = st.file_uploader(
        "Upload Merchant Data (CSV/XLS/XLSX)", type=["csv", "xlsx", "xls"], key="data_file", label_visibility="visible"
    )
    if not uploaded:
        st.info("📁 Please upload a CSV, XLS, or XLSX file to proceed.")
        return

    try:
        name = uploaded.name.lower()
        if name.endswith('.csv'):
            df = pd.read_csv(uploaded)
        else:
            df = pd.read_excel(uploaded)
    except Exception as e:
        st.error(f"Could not read file: {e}")
        return

    # Normalize column names
    df.columns = df.columns.str.strip()

    # Data enrichment
    with st.expander("Data Enrichment", expanded=True):
        df['Enriched_Location'] = df.get('Country', np.nan)
        df['Enriched_RiskScore'] = np.random.randint(1, 100, len(df))
        st.success("Data enrichment applied")
        st.metric("Rows enriched", len(df))
        st.metric("Avg Risk Score", int(df['Enriched_RiskScore'].mean()))
        st.dataframe(df.head(10))

    # Run screening & fraud
    df, summary_df, fraud_logic_cols = run_fraud_analysis(df)

    # Screening summary
    with st.expander("Sanctions/Screening", expanded=True):
        checks = pd.DataFrame({
            'Check': ['Sanction', 'HighRiskCountry', 'HighRiskCategory', 'HighAmount', 'RecentRegistration'],
            'Hits': [
                int(df['SanctionedFlag'].sum()),
                int(df['HighRiskCountryFlag'].sum()),
                int(df['HighRiskCategoryFlag'].sum()),
                int(df['HighAmountFlag'].sum()),
                int(df['RecentRegistrationFlag'].sum()),
            ]
        })
        st.table(checks)
        st.bar_chart(checks.set_index('Check'))

    # Executive summary export
    with pd.ExcelWriter(BytesIO(), engine='xlsxwriter') as writer:
        summary_df.to_excel(writer, sheet_name='Executive Summary', index=False)
        writer.book.close()
        excel_io = writer.path
    st.download_button(
        "Download Excel Report",
        data=excel_io.read(),
        file_name="executive_summary.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )

    # PowerPoint report
    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Executive Summary"
    summary_text = (
        f"Total merchants: {len(df)}\n"
        f"Sanction matches: {int(df['SanctionedFlag'].sum())}\n"
        f"High-risk jurisdictions: {int(df['HighRiskCountryFlag'].sum())}\n"
        f"Suspicious merchants: {int(summary_df['Count'].iloc[3])}"
    )
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = summary_text

    # Chart slide
    try:
        chart_layout = prs.slide_layouts[5]  # Title only layout
    except IndexError:
        chart_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(chart_layout)
    slide.shapes.title.text = "Fraud Logic Flags Distribution"

    buf = BytesIO()
    fig, ax = plt.subplots(figsize=(10, 6))
    pd.Series(df[fraud_logic_cols].sum()).sort_values(ascending=False).plot(kind='bar', ax=ax)
    ax.set_ylabel('Flag Count')
    ax.set_title('Fraud Flag Distribution')
    plt.tight_layout()
    fig.savefig(buf, format='png', dpi=120)
    plt.close(fig)
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(0.5), Inches(1.5), width=Inches(9))

    # Save PPT
    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    st.download_button(
        "Download PowerPoint Report", 
        data=pptx_io.read(), 
        file_name="executive_summary.pptx", 
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

    if st.button("Back to Main"):
        st.session_state.page = 'splash'

# --- Run pages ---
if st.session_state.page == 'splash':
    splash_page()
else:
    analysis_page()
